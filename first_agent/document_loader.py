import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import markdown
from bs4 import BeautifulSoup  # 用于清理 HTML

def load_text_from_pdf(filepath: str) -> str:
    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def load_text_from_docx(filepath: str) -> str:
    doc = Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def load_text_from_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_text_from_markdown(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8") as f:
        md_text = f.read()
    # 转换为 HTML 再提取纯文本（可选）
    html = markdown.markdown(md_text)
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text()

def load_document(filepath: str) -> str:
    """根据扩展名自动选择合适的解析器"""
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return load_text_from_pdf(filepath)
    elif ext == ".docx":
        return load_text_from_docx(filepath)
    elif ext == ".pptx":
        return load_text_from_pptx(filepath)
    elif ext == ".md":
        return load_text_from_markdown(filepath)
    elif ext == ".txt":
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    else:
        raise ValueError(f"不支持的文件类型: {ext}")